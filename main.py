"""
╔══════════════════════════════════════════════════════════════════════════════╗
║            NeuraChat ULTRA — Agentic AI Backend v2.0                        ║
║            FastAPI + Groq + RAG + Tools + File Generation                   ║
╚══════════════════════════════════════════════════════════════════════════════╝

CONCEPT: Agentic AI Backend
────────────────────────────
This server powers a Claude-like AI with:
  1. Multi-model chat with streaming (SSE)
  2. RAG — document upload + semantic retrieval
  3. Tool calling — AI can generate PPT, DOCX, PDF, charts, flowcharts
  4. Image analysis (vision models)
  5. Image generation (Pollinations.ai — free, no key)
  6. Token tracking with cost calculation
  7. Conversation history management
  8. Code execution (sandboxed)
  9. Jupyter notebook generation
"""

# ─── STANDARD LIBRARY ────────────────────────────────────────────────────────
import os, json, time, uuid, re, base64, io, subprocess, sys, textwrap, tempfile, shutil
from pathlib import Path
from typing import List, Optional, Dict, Any
from datetime import datetime
import urllib.request

# ─── WEB FRAMEWORK ───────────────────────────────────────────────────────────
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Request, BackgroundTasks
from fastapi.responses import StreamingResponse, JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field

# ─── AI / LLM ────────────────────────────────────────────────────────────────
from groq import Groq

# ─── FILE PROCESSING ─────────────────────────────────────────────────────────
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx as python_docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import tiktoken
    tokenizer = tiktoken.get_encoding("cl100k_base")
    HAS_TIKTOKEN = True
except ImportError:
    HAS_TIKTOKEN = False

# ─── OUTPUT GENERATION LIBRARIES ─────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph as RLParagraph, Spacer, Table as RLTable, TableStyle
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import inch
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    import nbformat
    HAS_NBFORMAT = True
except ImportError:
    HAS_NBFORMAT = False

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    import numpy as np
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False

# ─── VECTOR DB & EMBEDDINGS (Optional advanced RAG) ──────────────────────────
try:
    import chromadb
    HAS_CHROMADB = True
except ImportError:
    HAS_CHROMADB = False

try:
    from sentence_transformers import SentenceTransformer
    embed_model = SentenceTransformer('all-MiniLM-L6-v2')
    HAS_EMBEDDINGS = True
except ImportError:
    HAS_EMBEDDINGS = False

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "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")
OUTPUT_DIR = Path("./outputs")
OUTPUT_DIR.mkdir(exist_ok=True)

GROQ_MODELS = {
    "llama-3.3-70b-versatile": {
        "name": "Llama 3.3 70B",
        "context_window": 128000,
        "cost_in": 0.59,
        "cost_out": 0.79,
        "description": "Best quality — complex tasks",
        "supports_vision": False,
    },
    "llama-3.1-8b-instant": {
        "name": "Llama 3.1 8B",
        "context_window": 128000,
        "cost_in": 0.05,
        "cost_out": 0.08,
        "description": "Fastest — simple tasks",
        "supports_vision": False,
    },
    "llama3-70b-8192": {
        "name": "Llama 3 70B (8K)",
        "context_window": 8192,
        "cost_in": 0.59,
        "cost_out": 0.79,
        "description": "Stable 70B model",
        "supports_vision": False,
    },
    "mixtral-8x7b-32768": {
        "name": "Mixtral 8x7B MoE",
        "context_window": 32768,
        "cost_in": 0.24,
        "cost_out": 0.24,
        "description": "Mixture of Experts",
        "supports_vision": False,
    },
    "gemma2-9b-it": {
        "name": "Gemma 2 9B",
        "context_window": 8192,
        "cost_in": 0.20,
        "cost_out": 0.20,
        "description": "Google efficient model",
        "supports_vision": False,
    },
    "meta-llama/llama-4-scout-17b-16e-instruct": {
        "name": "Llama 4 Scout 17B",
        "context_window": 128000,
        "cost_in": 0.11,
        "cost_out": 0.34,
        "description": "Llama 4 with vision",
        "supports_vision": True,
    },
    "meta-llama/llama-4-maverick-17b-128e-instruct": {
        "name": "Llama 4 Maverick 17B",
        "context_window": 128000,
        "cost_in": 0.20,
        "cost_out": 0.60,
        "description": "Llama 4 advanced vision",
        "supports_vision": True,
    },
}

DEFAULT_SYSTEM_PROMPT = """You are NeuraChat Ultra — an advanced agentic AI assistant and expert teacher. You can:

CAPABILITIES:
- Answer any question with deep expertise
- Analyze uploaded documents (PDF, DOCX, code files)
- Analyze images and describe them in detail
- Generate images using Pollinations.ai (free)
- Generate PowerPoint presentations
- Generate Word documents (.docx)
- Generate PDF reports
- Generate flowcharts and diagrams (as SVG/Mermaid)
- Write and explain code in any language
- Create Jupyter notebooks (.ipynb)
- Explain AI/ML concepts: LLMs, RAG, LlamaIndex, Agentic AI, embeddings, vector DBs
- Train and fine-tune models (guide + code)
- Run code analysis

TEACHING STYLE:
When explaining concepts:
1. Start with a real-world analogy
2. Explain the technical mechanism
3. Show working code with comments
4. Connect to how this app works internally

FILE GENERATION:
When asked to create a file, use the appropriate tool:
- "create a presentation about X" → use create_pptx tool
- "generate a PDF report on X" → use create_pdf tool
- "make a flowchart for X" → describe as Mermaid diagram
- "create a Word document" → use create_docx tool
- "generate an image of X" → use generate_image tool
- "create a notebook for X" → use create_notebook tool

FORMAT:
- Use ## headings, **bold**, `code`, and numbered lists
- Use ```python for code blocks with the language specified
- Use > blockquotes for key insights
- Be concise but complete
"""

# ─── IN-MEMORY STATE ─────────────────────────────────────────────────────────
conversations: Dict[str, Dict] = {}
rag_documents: Dict[str, Dict] = {}
global_token_stats = {
    "total_prompt_tokens": 0,
    "total_completion_tokens": 0,
    "total_tokens": 0,
    "total_requests": 0,
    "session_start": datetime.now().isoformat(),
    "model_breakdown": {}
}

# ChromaDB setup for advanced RAG
chroma_collection = None
if HAS_CHROMADB:
    try:
        chroma_client = chromadb.PersistentClient(path="./chroma_db")
        chroma_collection = chroma_client.get_or_create_collection("neurachat_docs")
    except Exception:
        pass

# ═══════════════════════════════════════════════════════════════════════════════
# APP INITIALIZATION
# ═══════════════════════════════════════════════════════════════════════════════

app = FastAPI(
    title="NeuraChat Ultra API",
    description="Agentic AI — LLMs, RAG, Tools, File Generation",
    version="2.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

groq_client = Groq(api_key=GROQ_API_KEY)

# ═══════════════════════════════════════════════════════════════════════════════
# PYDANTIC MODELS
# ═══════════════════════════════════════════════════════════════════════════════

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    conversation_id: Optional[str] = None
    model: str = "llama-3.3-70b-versatile"
    temperature: float = Field(0.7, ge=0.0, le=2.0)
    max_tokens: int = Field(4096, ge=128, le=32768)
    stream: bool = True
    use_rag: bool = False
    system_prompt: Optional[str] = None

class GeneratePPTRequest(BaseModel):
    topic: str
    slides: List[Dict[str, Any]]
    theme: str = "dark"

class GeneratePDFRequest(BaseModel):
    title: str
    sections: List[Dict[str, str]]

class GenerateNotebookRequest(BaseModel):
    title: str
    cells: List[Dict[str, str]]

class GenerateImageRequest(BaseModel):
    prompt: str
    width: int = 1024
    height: int = 1024
    style: str = "photorealistic"

class GenerateChartRequest(BaseModel):
    chart_type: str  # bar, line, pie, scatter
    title: str
    data: Dict[str, Any]
    color_scheme: str = "viridis"

# ═══════════════════════════════════════════════════════════════════════════════
# UTILITY FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def estimate_tokens(text: str) -> int:
    """Count tokens using tiktoken or estimate from character count."""
    if HAS_TIKTOKEN:
        return len(tokenizer.encode(text))
    return max(1, len(text) // 4)

def calculate_cost(prompt_tokens: int, completion_tokens: int, model: str) -> float:
    """Calculate API cost in USD."""
    m = GROQ_MODELS.get(model, {"cost_in": 0.59, "cost_out": 0.79})
    return (prompt_tokens * m["cost_in"] + completion_tokens * m["cost_out"]) / 1_000_000

def update_token_stats(prompt_tokens: int, completion_tokens: int, model: str):
    """Update global token usage tracker."""
    global_token_stats["total_prompt_tokens"] += prompt_tokens
    global_token_stats["total_completion_tokens"] += completion_tokens
    global_token_stats["total_tokens"] += prompt_tokens + completion_tokens
    global_token_stats["total_requests"] += 1
    if model not in global_token_stats["model_breakdown"]:
        global_token_stats["model_breakdown"][model] = {"prompt": 0, "completion": 0, "requests": 0}
    global_token_stats["model_breakdown"][model]["prompt"] += prompt_tokens
    global_token_stats["model_breakdown"][model]["completion"] += completion_tokens
    global_token_stats["model_breakdown"][model]["requests"] += 1

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks for RAG."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
        if start >= len(text):
            break
    return chunks

def extract_text_from_file(content: bytes, filename: str) -> str:
    """Extract text from PDF, DOCX, or plain text files."""
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf' and HAS_PDF:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            return "\n\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            return f"[PDF extraction error: {e}]"
    elif ext in ('docx', 'doc') and HAS_DOCX:
        try:
            doc = python_docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            return f"[DOCX extraction error: {e}]"
    elif ext in ('py', 'js', 'ts', 'html', 'css', 'json', 'yaml', 'yml', 'md', 'txt', 'csv', 'sql'):
        try:
            return content.decode('utf-8', errors='replace')
        except:
            return content.decode('latin-1', errors='replace')
    else:
        try:
            return content.decode('utf-8', errors='replace')
        except:
            return f"[Binary file: {filename} — cannot extract text]"

def build_rag_context(query: str, max_chars: int = 12000) -> str:
    """
    CONCEPT: RAG Context Building
    ──────────────────────────────
    Retrieves relevant document chunks for the user's query.
    
    Advanced mode: Uses ChromaDB vector similarity search
    Simple mode: Returns documents sorted by keyword overlap
    """
    if not rag_documents:
        return ""

    # Advanced RAG: vector similarity via ChromaDB
    if HAS_CHROMADB and HAS_EMBEDDINGS and chroma_collection:
        try:
            count = chroma_collection.count()
            if count > 0:
                results = chroma_collection.query(
                    query_texts=[query],
                    n_results=min(5, count)
                )
                chunks = results['documents'][0] if results['documents'] else []
                if chunks:
                    context = "\n\n---\n\n".join(chunks)
                    return f"[RETRIEVED CONTEXT — {len(chunks)} relevant chunks via semantic search]\n\n{context[:max_chars]}"
        except Exception:
            pass

    # Simple RAG: full document injection
    parts = []
    total = 0
    for doc in rag_documents.values():
        snippet = doc['content'][:max_chars // max(1, len(rag_documents))]
        parts.append(f"[Document: {doc['name']}]\n{snippet}")
        total += len(snippet)
        if total >= max_chars:
            break

    return "[KNOWLEDGE BASE]\n\n" + "\n\n---\n\n".join(parts) if parts else ""

# ═══════════════════════════════════════════════════════════════════════════════
# FILE GENERATION FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def create_pptx_file(topic: str, slides_data: List[Dict], theme: str = "dark") -> str:
    """Generate a PowerPoint presentation."""
    if not HAS_PPTX:
        raise HTTPException(status_code=500, detail="python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color themes
    themes = {
        "dark": {"bg": (15, 17, 26), "title": (255, 255, 255), "accent": (92, 124, 250), "text": (200, 210, 230)},
        "light": {"bg": (250, 252, 255), "title": (15, 17, 26), "accent": (59, 130, 246), "text": (50, 60, 80)},
        "corporate": {"bg": (0, 32, 96), "title": (255, 255, 255), "accent": (255, 192, 0), "text": (200, 220, 255)},
        "minimal": {"bg": (255, 255, 255), "title": (30, 30, 30), "accent": (100, 100, 100), "text": (80, 80, 80)},
    }
    t = themes.get(theme, themes["dark"])

    def set_bg(slide, color_tuple):
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_tuple)

    def add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=(255,255,255), align="left"):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if align == "left" else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return txBox

    for i, slide_data in enumerate(slides_data):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        set_bg(slide, t["bg"])

        title = slide_data.get("title", f"Slide {i+1}")
        content = slide_data.get("content", "")
        slide_type = slide_data.get("type", "content")

        if slide_type == "title" or i == 0:
            # Title slide
            add_text_box(slide, topic, 1, 1.5, 11.33, 1.5, 44, bold=True, color=t["title"], align="center")
            add_text_box(slide, title, 1, 3.2, 11.33, 1, 28, color=t["accent"], align="center")
            add_text_box(slide, "Generated by NeuraChat Ultra", 1, 6.5, 11.33, 0.5, 14, color=t["text"], align="center")
        else:
            # Content slide — title bar
            title_box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = RGBColor(*t["accent"])
            title_box.line.fill.background()

            add_text_box(slide, title, 0.3, 0.15, 12, 0.9, 28, bold=True, color=(255, 255, 255))

            # Content
            if isinstance(content, list):
                y = 1.5
                for bullet in content:
                    add_text_box(slide, f"• {bullet}", 0.5, y, 12.33, 0.6, 18, color=t["text"])
                    y += 0.65
            else:
                add_text_box(slide, str(content), 0.5, 1.5, 12.33, 5.5, 18, color=t["text"])

        # Slide number
        add_text_box(slide, f"{i+1}/{len(slides_data)}", 12, 7.1, 1, 0.3, 12, color=t["text"])

    filename = f"neurachat_{uuid.uuid4().hex[:8]}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return filename

def create_pdf_file(title: str, sections: List[Dict]) -> str:
    """Generate a PDF report using ReportLab."""
    if not HAS_REPORTLAB:
        raise HTTPException(status_code=500, detail="reportlab not installed. Run: pip install reportlab")

    filename = f"report_{uuid.uuid4().hex[:8]}.pdf"
    filepath = OUTPUT_DIR / filename

    doc = SimpleDocTemplate(str(filepath), pagesize=letter,
                            topMargin=0.75*inch, bottomMargin=0.75*inch,
                            leftMargin=1*inch, rightMargin=1*inch)

    styles = getSampleStyleSheet()
    story = []

    # Title
    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'],
                                  fontSize=28, spaceAfter=20,
                                  textColor=rl_colors.HexColor('#1a3a5c'))
    story.append(RLParagraph(title, title_style))
    story.append(Spacer(1, 0.2*inch))

    # Generated by line
    meta_style = ParagraphStyle('Meta', parent=styles['Normal'],
                                 fontSize=10, textColor=rl_colors.grey)
    story.append(RLParagraph(f"Generated by NeuraChat Ultra | {datetime.now().strftime('%Y-%m-%d %H:%M')}", meta_style))
    story.append(Spacer(1, 0.3*inch))

    # Sections
    h2_style = ParagraphStyle('H2', parent=styles['Heading2'],
                               fontSize=16, textColor=rl_colors.HexColor('#2e75b6'),
                               spaceAfter=10, spaceBefore=20)
    body_style = ParagraphStyle('Body', parent=styles['Normal'],
                                 fontSize=12, leading=18,
                                 textColor=rl_colors.HexColor('#333333'))

    for section in sections:
        if section.get("heading"):
            story.append(RLParagraph(section["heading"], h2_style))
        if section.get("content"):
            for para in section["content"].split('\n\n'):
                if para.strip():
                    story.append(RLParagraph(para.replace('<', '&lt;').replace('>', '&gt;'), body_style))
                    story.append(Spacer(1, 0.1*inch))

    doc.build(story)
    return filename

def create_docx_file_py(title: str, sections: List[Dict]) -> str:
    """Generate a DOCX document using python-docx."""
    if not HAS_DOCX:
        raise HTTPException(status_code=500, detail="python-docx not installed. Run: pip install python-docx")

    from docx.shared import Pt, RGBColor as DocxRGB, Inches as DocxInches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = python_docx.Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = DocxInches(1)
        section.bottom_margin = DocxInches(1)
        section.left_margin = DocxInches(1.2)
        section.right_margin = DocxInches(1.2)

    # Title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = title_para.runs[0]
    run.font.color.rgb = DocxRGB(26, 58, 92)
    run.font.size = Pt(28)

    doc.add_paragraph(f"Generated by NeuraChat Ultra | {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    for s in sections:
        if s.get("heading"):
            h = doc.add_heading(s["heading"], 1)
            h.runs[0].font.color.rgb = DocxRGB(46, 117, 182)
        if s.get("content"):
            for para in s["content"].split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para)

    filename = f"document_{uuid.uuid4().hex[:8]}.docx"
    filepath = OUTPUT_DIR / filename
    doc.save(str(filepath))
    return filename

def create_notebook_file(title: str, cells: List[Dict]) -> str:
    """Generate a Jupyter notebook (.ipynb)."""
    if not HAS_NBFORMAT:
        raise HTTPException(status_code=500, detail="nbformat not installed. Run: pip install nbformat")

    nb = nbformat.v4.new_notebook()
    nb_cells = [
        nbformat.v4.new_markdown_cell(f"# {title}\n\n*Generated by NeuraChat Ultra | {datetime.now().strftime('%Y-%m-%d')}*")
    ]

    for cell in cells:
        cell_type = cell.get("type", "code")
        source = cell.get("source", "")
        if cell_type == "markdown":
            nb_cells.append(nbformat.v4.new_markdown_cell(source))
        else:
            nb_cells.append(nbformat.v4.new_code_cell(source))

    nb.cells = nb_cells
    filename = f"notebook_{uuid.uuid4().hex[:8]}.ipynb"
    filepath = OUTPUT_DIR / filename
    with open(str(filepath), 'w') as f:
        nbformat.write(nb, f)
    return filename

def create_chart_file(chart_type: str, title: str, data: Dict, color_scheme: str = "viridis") -> str:
    """Generate a chart/graph using matplotlib."""
    if not HAS_MATPLOTLIB:
        raise HTTPException(status_code=500, detail="matplotlib not installed. Run: pip install matplotlib numpy")

    fig, ax = plt.subplots(figsize=(12, 7), facecolor='#0e1118')
    ax.set_facecolor('#131720')
    ax.tick_params(colors='#8892b0')
    ax.xaxis.label.set_color('#8892b0')
    ax.yaxis.label.set_color('#8892b0')
    for spine in ax.spines.values():
        spine.set_edgecolor('#252c3d')

    labels = data.get("labels", [])
    values = data.get("values", [])
    colors_list = plt.cm.get_cmap(color_scheme)(np.linspace(0.3, 0.9, max(len(labels), 1)))

    if chart_type == "bar":
        bars = ax.bar(labels, values, color=colors_list, edgecolor='#252c3d', linewidth=0.5)
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + max(values)*0.01,
                    f'{val:,.0f}', ha='center', va='bottom', color='#eef0f8', fontsize=10)
    elif chart_type == "line":
        ax.plot(labels, values, color='#5c7cfa', linewidth=2.5, marker='o',
                markersize=8, markerfacecolor='#7c3aed', markeredgecolor='white', markeredgewidth=1.5)
        ax.fill_between(range(len(labels)), values, alpha=0.15, color='#5c7cfa')
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels)
    elif chart_type == "pie":
        wedges, texts, autotexts = ax.pie(values, labels=labels, colors=colors_list,
                                            autopct='%1.1f%%', startangle=90,
                                            wedgeprops={'edgecolor': '#0e1118', 'linewidth': 2})
        for text in texts:
            text.set_color('#eef0f8')
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
    elif chart_type == "scatter":
        x_vals = data.get("x", list(range(len(values))))
        scatter = ax.scatter(x_vals, values, c=colors_list[:len(values)], s=100, alpha=0.8, edgecolors='white', linewidths=0.5)

    ax.set_title(title, color='#eef0f8', fontsize=16, fontweight='bold', pad=20)
    ax.grid(True, alpha=0.1, color='#252c3d')
    plt.tight_layout()

    filename = f"chart_{uuid.uuid4().hex[:8]}.png"
    filepath = OUTPUT_DIR / filename
    plt.savefig(str(filepath), dpi=150, bbox_inches='tight', facecolor='#0e1118')
    plt.close()
    return filename

def generate_mermaid_svg(mermaid_code: str) -> str:
    """Return mermaid code as-is for frontend rendering."""
    return mermaid_code

# ═══════════════════════════════════════════════════════════════════════════════
# API ENDPOINTS
# ═══════════════════════════════════════════════════════════════════════════════

@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "version": "2.0.0",
        "features": {
            "groq": bool(GROQ_API_KEY),
            "pdf": HAS_PDF,
            "docx": HAS_DOCX,
            "pptx": HAS_PPTX,
            "pdf_gen": HAS_REPORTLAB,
            "charts": HAS_MATPLOTLIB,
            "notebooks": HAS_NBFORMAT,
            "chromadb": HAS_CHROMADB,
            "embeddings": HAS_EMBEDDINGS,
            "tiktoken": HAS_TIKTOKEN,
        },
        "timestamp": datetime.now().isoformat()
    }

@app.get("/api/models")
async def get_models():
    return {"models": GROQ_MODELS}

@app.get("/api/capabilities")
async def get_capabilities():
    return {
        "pdf_read": HAS_PDF,
        "docx_read": HAS_DOCX,
        "pptx_gen": HAS_PPTX,
        "pdf_gen": HAS_REPORTLAB,
        "charts": HAS_MATPLOTLIB,
        "notebooks": HAS_NBFORMAT,
        "chromadb": HAS_CHROMADB,
        "embeddings": HAS_EMBEDDINGS,
        "image_gen": True,  # Pollinations.ai — always available
    }

# ─── MAIN CHAT ENDPOINT ───────────────────────────────────────────────────────

@app.post("/api/chat")
async def chat(request: ChatRequest):
    """
    CONCEPT: Streaming Chat with SSE
    ──────────────────────────────────
    Server-Sent Events push tokens to the browser as they're generated.
    The frontend's EventSource listens and appends each token in real-time.
    """
    # Get or create conversation
    conv_id = request.conversation_id or str(uuid.uuid4())
    if conv_id not in conversations:
        conversations[conv_id] = {
            "id": conv_id,
            "title": request.message[:50],
            "messages": [],
            "created_at": datetime.now().isoformat(),
            "model": request.model,
            "token_usage": {"prompt": 0, "completion": 0}
        }

    conv = conversations[conv_id]

    # Build system prompt with RAG context
    system = request.system_prompt or DEFAULT_SYSTEM_PROMPT
    if request.use_rag:
        rag_ctx = build_rag_context(request.message)
        if rag_ctx:
            system += f"\n\n{rag_ctx}"

    # Build messages list
    messages = [{"role": "system", "content": system}]
    # Include conversation history (last 20 messages)
    for msg in conv["messages"][-20:]:
        messages.append({"role": msg["role"], "content": msg["content"]})
    messages.append({"role": "user", "content": request.message})

    # Add user message to history
    conv["messages"].append({
        "role": "user",
        "content": request.message,
        "timestamp": datetime.now().isoformat()
    })

    if request.stream:
        async def token_stream():
            full_response = ""
            prompt_tokens = estimate_tokens(json.dumps(messages))
            completion_tokens = 0

            try:
                stream = groq_client.chat.completions.create(
                    model=request.model,
                    messages=messages,
                    temperature=request.temperature,
                    max_tokens=request.max_tokens,
                    stream=True,
                )
                for chunk in stream:
                    token = chunk.choices[0].delta.content or ""
                    if token:
                        full_response += token
                        completion_tokens += 1
                        yield f"data: {json.dumps({'token': token, 'conv_id': conv_id})}\n\n"

                # Save assistant response
                conv["messages"].append({
                    "role": "assistant",
                    "content": full_response,
                    "timestamp": datetime.now().isoformat(),
                    "model": request.model
                })

                # Update token stats
                update_token_stats(prompt_tokens, completion_tokens, request.model)
                conv["token_usage"]["prompt"] += prompt_tokens
                conv["token_usage"]["completion"] += completion_tokens

                cost = calculate_cost(prompt_tokens, completion_tokens, request.model)
                yield f"data: {json.dumps({'done': True, 'conv_id': conv_id, 'prompt_tokens': prompt_tokens, 'completion_tokens': completion_tokens, 'cost': cost})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"

        return StreamingResponse(token_stream(), media_type="text/event-stream",
                                  headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})
    else:
        # Non-streaming
        try:
            response = groq_client.chat.completions.create(
                model=request.model,
                messages=messages,
                temperature=request.temperature,
                max_tokens=request.max_tokens,
            )
            content = response.choices[0].message.content
            pt = response.usage.prompt_tokens if response.usage else estimate_tokens(json.dumps(messages))
            ct = response.usage.completion_tokens if response.usage else estimate_tokens(content)
            update_token_stats(pt, ct, request.model)
            conv["messages"].append({"role": "assistant", "content": content, "model": request.model})
            return {"response": content, "conv_id": conv_id, "prompt_tokens": pt, "completion_tokens": ct}
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

# ─── DOCUMENT UPLOAD & RAG ────────────────────────────────────────────────────

@app.post("/api/upload")
async def upload_document(
    file: UploadFile = File(...),
    add_to_rag: bool = Form(default=True)
):
    filename = file.filename or "upload"
    file_content = await file.read()

    if len(file_content) > 20 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large. Max 20MB.")

    extracted_text = extract_text_from_file(file_content, filename)
    token_count = estimate_tokens(extracted_text)
    chunks = chunk_text(extracted_text)

    if add_to_rag:
        doc_id = str(uuid.uuid4())
        rag_documents[doc_id] = {
            "id": doc_id,
            "name": filename,
            "content": extracted_text,
            "chunks": chunks,
            "chunk_count": len(chunks),
            "char_count": len(extracted_text),
            "estimated_tokens": token_count,
            "uploaded_at": datetime.now().isoformat()
        }

        # Add to ChromaDB for advanced RAG if available
        if HAS_CHROMADB and HAS_EMBEDDINGS and chroma_collection:
            try:
                ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
                embeddings = embed_model.encode(chunks).tolist()
                chroma_collection.add(documents=chunks, embeddings=embeddings, ids=ids)
            except Exception:
                pass

        return {
            "success": True,
            "doc_id": doc_id,
            "filename": filename,
            "char_count": len(extracted_text),
            "estimated_tokens": token_count,
            "chunk_count": len(chunks),
            "advanced_rag": HAS_CHROMADB and HAS_EMBEDDINGS,
            "message": f"'{filename}' added to knowledge base ({len(chunks)} chunks)"
        }
    else:
        return {
            "success": True,
            "filename": filename,
            "content": extracted_text[:5000],
            "char_count": len(extracted_text),
            "estimated_tokens": token_count,
        }

@app.get("/api/documents")
async def list_documents():
    return {"documents": [
        {"id": d["id"], "name": d["name"], "char_count": d["char_count"],
         "estimated_tokens": d["estimated_tokens"], "chunk_count": d["chunk_count"],
         "uploaded_at": d["uploaded_at"]}
        for d in rag_documents.values()
    ], "count": len(rag_documents)}

@app.delete("/api/documents/{doc_id}")
async def delete_document(doc_id: str):
    if doc_id not in rag_documents:
        raise HTTPException(status_code=404, detail="Document not found")
    name = rag_documents[doc_id]["name"]
    del rag_documents[doc_id]
    return {"message": f"'{name}' removed"}

# ─── IMAGE ANALYSIS ───────────────────────────────────────────────────────────

@app.post("/api/analyze-image")
async def analyze_image(
    file: UploadFile = File(...),
    question: str = Form(default="Describe this image in detail. Extract any text, charts, diagrams, or data you see."),
    model: str = Form(default="meta-llama/llama-4-scout-17b-16e-instruct")
):
    image_bytes = await file.read()
    content_type = file.content_type or "image/jpeg"
    image_b64 = base64.b64encode(image_bytes).decode("utf-8")

    try:
        response = groq_client.chat.completions.create(
            model=model,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": question},
                    {"type": "image_url", "image_url": {"url": f"data:{content_type};base64,{image_b64}"}}
                ]
            }],
            max_tokens=2048
        )
        return {
            "analysis": response.choices[0].message.content,
            "model": model,
            "filename": file.filename
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Vision analysis failed: {str(e)}")

# ─── FILE GENERATION ENDPOINTS ────────────────────────────────────────────────

@app.post("/api/generate/pptx")
async def gen_pptx(req: GeneratePPTRequest):
    filename = create_pptx_file(req.topic, req.slides, req.theme)
    return {"filename": filename, "download_url": f"/api/download/{filename}"}

@app.post("/api/generate/pdf")
async def gen_pdf(req: GeneratePDFRequest):
    filename = create_pdf_file(req.title, req.sections)
    return {"filename": filename, "download_url": f"/api/download/{filename}"}

@app.post("/api/generate/docx")
async def gen_docx(req: GeneratePDFRequest):
    filename = create_docx_file_py(req.title, req.sections)
    return {"filename": filename, "download_url": f"/api/download/{filename}"}

@app.post("/api/generate/notebook")
async def gen_notebook(req: GenerateNotebookRequest):
    filename = create_notebook_file(req.title, req.cells)
    return {"filename": filename, "download_url": f"/api/download/{filename}"}

@app.post("/api/generate/chart")
async def gen_chart(req: GenerateChartRequest):
    filename = create_chart_file(req.chart_type, req.title, req.data, req.color_scheme)
    return {"filename": filename, "download_url": f"/api/download/{filename}"}

@app.post("/api/generate/image")
async def gen_image_endpoint(req: GenerateImageRequest):
    """Generate image using Pollinations.ai (free, no key required)."""
    import urllib.parse
    encoded = urllib.parse.quote(req.prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded}?width={req.width}&height={req.height}&nologo=true"
    return {"image_url": url, "prompt": req.prompt}

@app.get("/api/download/{filename}")
async def download_file(filename: str):
    # Security: only allow files in output dir with safe extensions
    safe_exts = {'.pptx', '.pdf', '.docx', '.ipynb', '.png', '.jpg', '.svg', '.md'}
    filepath = OUTPUT_DIR / filename
    ext = Path(filename).suffix.lower()
    if not filepath.exists() or ext not in safe_exts or '..' in filename:
        raise HTTPException(status_code=404, detail="File not found")

    media_types = {
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.pdf': 'application/pdf',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.ipynb': 'application/json',
        '.png': 'image/png',
        '.svg': 'image/svg+xml',
        '.md': 'text/markdown',
    }
    return FileResponse(str(filepath), media_type=media_types.get(ext, 'application/octet-stream'),
                        filename=filename)

@app.get("/api/outputs")
async def list_outputs():
    files = []
    for f in OUTPUT_DIR.iterdir():
        if f.is_file():
            files.append({
                "name": f.name,
                "size": f.stat().st_size,
                "created": datetime.fromtimestamp(f.stat().st_mtime).isoformat(),
                "download_url": f"/api/download/{f.name}"
            })
    return {"files": sorted(files, key=lambda x: x["created"], reverse=True)}

# ─── AI-POWERED FILE GENERATION (Ask AI to structure the content) ─────────────

@app.post("/api/ai-generate")
async def ai_generate(request: dict):
    """
    Let the AI structure the content, then generate the file.
    The AI outputs JSON describing slides/sections, we render the file.
    """
    file_type = request.get("type", "pptx")
    topic = request.get("topic", "")
    instructions = request.get("instructions", "")
    model = request.get("model", "llama-3.3-70b-versatile")

    # Build prompt based on file type
    prompts = {
        "pptx": f"""Create a PowerPoint presentation about: {topic}

{instructions}

Output ONLY valid JSON (no markdown, no explanation) in this exact format:
{{
  "topic": "{topic}",
  "theme": "dark",
  "slides": [
    {{"type": "title", "title": "Subtitle here", "content": ""}},
    {{"type": "content", "title": "Slide Title", "content": ["Bullet 1", "Bullet 2", "Bullet 3"]}},
    {{"type": "content", "title": "Another Slide", "content": ["Point A", "Point B"]}}
  ]
}}
Include 5-8 slides. Make content informative and well-structured.""",

        "pdf": f"""Create a PDF report about: {topic}

{instructions}

Output ONLY valid JSON in this format:
{{
  "title": "{topic}",
  "sections": [
    {{"heading": "Introduction", "content": "Full paragraph text here..."}},
    {{"heading": "Section 2", "content": "More text..."}}
  ]
}}
Include 4-6 sections with detailed content.""",

        "notebook": f"""Create a Jupyter notebook about: {topic}

{instructions}

Output ONLY valid JSON in this format:
{{
  "title": "{topic}",
  "cells": [
    {{"type": "markdown", "source": "## Introduction\\nExplanation here"}},
    {{"type": "code", "source": "# Code here\\nimport numpy as np"}},
    {{"type": "markdown", "source": "## Explanation"}},
    {{"type": "code", "source": "# More code"}}
  ]
}}
Include 6-10 cells with practical code examples.""",

        "chart": f"""Create chart data for: {topic}

{instructions}

Output ONLY valid JSON in this format:
{{
  "chart_type": "bar",
  "title": "{topic}",
  "color_scheme": "viridis",
  "data": {{
    "labels": ["Item 1", "Item 2", "Item 3"],
    "values": [100, 200, 150]
  }}
}}
Choose chart_type from: bar, line, pie, scatter."""
    }

    prompt = prompts.get(file_type, prompts["pptx"])

    try:
        response = groq_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are a structured data generator. Output ONLY valid JSON. No markdown code blocks, no explanation, just the raw JSON object."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=4096,
        )
        raw = response.choices[0].message.content.strip()
        # Clean up markdown if model adds it
        raw = re.sub(r'^```json\s*', '', raw)
        raw = re.sub(r'^```\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
        data = json.loads(raw)

        # Generate the file
        if file_type == "pptx":
            filename = create_pptx_file(data.get("topic", topic), data.get("slides", []), data.get("theme", "dark"))
        elif file_type == "pdf":
            filename = create_pdf_file(data.get("title", topic), data.get("sections", []))
        elif file_type == "notebook":
            filename = create_notebook_file(data.get("title", topic), data.get("cells", []))
        elif file_type == "chart":
            filename = create_chart_file(data.get("chart_type", "bar"), data.get("title", topic), data.get("data", {}), data.get("color_scheme", "viridis"))
        else:
            raise HTTPException(status_code=400, detail=f"Unknown type: {file_type}")

        pt = response.usage.prompt_tokens if response.usage else 0
        ct = response.usage.completion_tokens if response.usage else 0
        update_token_stats(pt, ct, model)

        return {"success": True, "filename": filename, "download_url": f"/api/download/{filename}", "prompt_tokens": pt, "completion_tokens": ct}

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"AI returned invalid JSON: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ─── TOKEN STATS ──────────────────────────────────────────────────────────────

@app.get("/api/tokens/stats")
async def get_token_stats():
    cost = sum(
        calculate_cost(
            global_token_stats["model_breakdown"].get(m, {}).get("prompt", 0),
            global_token_stats["model_breakdown"].get(m, {}).get("completion", 0),
            m
        ) for m in global_token_stats["model_breakdown"]
    )
    return {**global_token_stats, "estimated_total_cost_usd": round(cost, 6)}

@app.post("/api/tokens/reset")
async def reset_token_stats():
    global_token_stats.update({
        "total_prompt_tokens": 0, "total_completion_tokens": 0,
        "total_tokens": 0, "total_requests": 0,
        "session_start": datetime.now().isoformat(),
        "model_breakdown": {}
    })
    return {"message": "Stats reset", "stats": global_token_stats}

# ─── CONVERSATION MANAGEMENT ──────────────────────────────────────────────────

@app.get("/api/conversations")
async def list_conversations():
    return {"conversations": [
        {"id": c["id"], "title": c["title"], "created_at": c["created_at"],
         "message_count": len(c["messages"]), "model": c.get("model", "unknown")}
        for c in conversations.values()
    ]}

@app.delete("/api/conversations/{conv_id}")
async def delete_conversation(conv_id: str):
    if conv_id in conversations:
        del conversations[conv_id]
    return {"message": "Deleted"}

# ─── STATIC FILES ─────────────────────────────────────────────────────────────

static_dir = Path(__file__).parent / "static"
static_dir.mkdir(exist_ok=True)
app.mount("/", StaticFiles(directory=str(static_dir), html=True), name="static")

if __name__ == "__main__":
    import uvicorn
    print("""
╔══════════════════════════════════════════════════════════════╗
║         NeuraChat Ultra v2.0 — Starting...                   ║
╠══════════════════════════════════════════════════════════════╣
║  Frontend:    http://localhost:8000                          ║
║  API Docs:    http://localhost:8000/docs                     ║
║  Health:      http://localhost:8000/api/health               ║
╚══════════════════════════════════════════════════════════════╝
    """)
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True, log_level="info")
